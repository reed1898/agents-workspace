from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLUE = RGBColor(20, 77, 163)
DARK = RGBColor(32, 33, 36)
GRAY = RGBColor(95, 99, 104)
LIGHT = RGBColor(245, 248, 255)


def set_bg(slide, color=RGBColor(255, 255, 255)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.95))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(31)
    r.font.bold = True
    r.font.color.rgb = DARK
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY


def add_bullets(slide, items, left=0.85, top=1.75, width=11.6, height=4.9):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for text, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        p.font.size = Pt(18 if level == 0 else 15)
        p.font.color.rgb = DARK if level == 0 else GRAY
        p.font.bold = (level == 0)
        p.space_after = Pt(6 if level == 0 else 3)


def add_footer(slide, page, name):
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.01))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(220, 225, 235); line.line.fill.background()
    f = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(10.5), Inches(0.3)).text_frame
    f.text = f"湖北厚爱健康管理有限公司｜{name}"
    f.paragraphs[0].font.size = Pt(10)
    f.paragraphs[0].font.color.rgb = GRAY
    pn = slide.shapes.add_textbox(Inches(12.2), Inches(6.86), Inches(0.5), Inches(0.3)).text_frame
    pn.text = str(page)
    pn.paragraphs[0].alignment = PP_ALIGN.RIGHT
    pn.paragraphs[0].font.size = Pt(10)
    pn.paragraphs[0].font.color.rgb = GRAY


def build_deck(slides, out_path, footer_name):
    prs = Presentation()
    for i, (title, subtitle, items) in enumerate(slides, 1):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s, LIGHT if i in [1, len(slides)] else RGBColor(255, 255, 255))
        add_title(s, title, subtitle)
        add_bullets(s, items)
        add_footer(s, i, footer_name)
    prs.save(out_path)


gov = [
    ("厚爱健康&生活大数据平台", "湖北省楚商联合会合作提案（政府汇报版）", [("医疗健康 + 生活消费 + 政企协同",0),("目标：稳民生、促消费、强产业、增税源",0)]),
    ("政策与现实背景", None, [("分级诊疗、医养结合、促消费政策持续推进",0),("民生服务与消费场景尚未形成闭环",0),("地方企业有供给，缺统一数字化交易与运营平台",0)]),
    ("项目定位", None, [("建设本地化‘健康+生活’大数据平台",0),("联动政府、联合会、企业、居民四方主体",0),("形成‘服务可及+消费可持续+数据可治理’体系",0)]),
    ("对民生治理价值", None, [("积分可用于水电煤/物业等民生支出兑换",0),("提升困难群体帮扶精准度与可持续性",0),("医疗健康服务下沉，增强基层可及性",0)]),
    ("对地方经济价值", None, [("本地消费留存提升，带动GDP与税收增长",0),("拉动本地商户销量与就业岗位",0),("形成区域消费数据资产，支撑政策优化",0)]),
    ("平台运行机制", None, [("消费获积分—积分兑换民生/商品—再消费",0),("构建持续激励而非一次性补贴",0),("沉淀居民健康与消费行为画像，支持精细治理",0)]),
    ("合规与风控设计", None, [("积分定位为消费激励工具，规则透明可审计",0),("数据分级授权、隐私保护、全链路审计",0),("交易结算隔离，商户准入与质量监管",0)]),
    ("试点建议（12个月）", None, [("0-3个月：1-2地市试点，打通核心流程",0),("3-6个月：扩展企业和场景，形成标准SOP",0),("6-12个月：省域复制，建立长期协同机制",0)]),
    ("合作诉求", None, [("建议成立政企联合推进专班",0),("开放民生场景试点与联合会企业组织协同",0),("共建‘湖北样板’，形成可复制模式",0)]),
    ("结语", None, [("以健康为入口，以消费为引擎，以数据为底座",0),("共建湖北本地经济内循环新范式",0)]),
]

biz = [
    ("厚爱健康&生活大数据平台", "湖北省楚商联合会合作提案（招商版）", [("帮助会员企业：低成本获客、持续复购、品牌增长",0)]),
    ("企业痛点", None, [("流量分散、获客成本高、复购弱",0),("中小企业数字化运营能力不足",0),("本地市场缺少可信联合品牌入口",0)]),
    ("平台能提供什么", None, [("楚商会员统一入驻与产品展示",0),("健康+生活高频流量入口导入",0),("积分激励提升下单率与复购率",0)]),
    ("增长模型", None, [("公域引流：联合会品牌活动/政企联动",0),("私域沉淀：会员体系+积分体系+内容运营",0),("复购驱动：民生兑换与健康服务场景绑定",0)]),
    ("商家收益", None, [("新增用户：获得本地稳定客源",0),("销售增长：客单、复购、转介绍提升",0),("品牌增信：楚商联合会背书提升转化",0)]),
    ("运营抓手", None, [("平台专场：楚商品牌周、健康惠民季",0),("用户运营：任务积分、会员等级、家庭账户",0),("数据运营：转化漏斗、品类热度、复购预警",0)]),
    ("入驻模式", None, [("标准入驻：上架+交易+结算",0),("联合营销：平台活动与品牌共投",0),("重点商家：定制增长方案与私域代运营",0)]),
    ("商业模式", None, [("平台服务费/技术服务费",0),("营销服务包与增值工具",0),("联合活动与品牌专区合作",0)]),
    ("阶段目标", None, [("3个月：100家标杆商户入驻",0),("6个月：形成稳定复购与会员增长曲线",0),("12个月：打造省内可复制商业样板",0)]),
    ("马上可做", None, [("确定首批行业（医药、家政、餐饮、零售）",0),("启动联合会专场招商会",0),("2周内完成首批商户上线",0)]),
]

roadshow = [
    ("厚爱健康&生活大数据平台", "综合路演版（政府+企业+生态伙伴）", [("一句话：把民生服务和消费增长做成同一个系统",0)]),
    ("我们看到的机会", None, [("民生刚需高频，消费激励可持续",0),("楚商组织力强，厚爱平台能力成熟",0),("现在是打造区域数字样板的窗口期",0)]),
    ("平台全景", None, [("一端连企业：产品与服务供给",0),("一端连居民：医疗健康与生活消费",0),("一端连政府：民生治理与经济监测",0)]),
    ("四大业务模块", None, [("医疗健康服务模块",0),("生活消费交易模块",0),("民生缴费兑换模块",0),("积分储备与养老激励模块（合规设计）",0)]),
    ("双边价值", None, [("左手企业：宣传+销售+复购",0),("右手政府：纾困+稳民生+促GDP税收",0),("中间居民：得实惠、得服务、得保障",0)]),
    ("核心飞轮", None, [("消费获积分→积分兑换民生/商品→提升消费意愿",0),("形成持续参与与长期留存",0),("数据反哺运营与政策优化",0)]),
    ("技术与数据底座", None, [("统一账户、统一积分、统一支付与结算",0),("多设备接入，支持健康数据采集",0),("可视化看板支持政府与企业决策",0)]),
    ("合规边界", None, [("积分为消费激励，不承诺金融收益",0),("数据合规：最小化采集、授权可追溯",0),("平台风控：反作弊、反刷单、反滥用",0)]),
    ("12个月落地计划", None, [("试点：打穿流程，验证模型",0),("扩面：增加企业与服务品类",0),("规模：省域复制，形成长期机制",0)]),
    ("合作机制", None, [("联合会牵头组织企业",0),("厚爱负责平台与运营",0),("政府指导民生与政策协同",0)]),
    ("预期成效", None, [("消费增长、税收留存、民生改善",0),("楚商品牌影响力提升",0),("形成湖北可复制的数字经济案例",0)]),
    ("谢谢", None, [("湖北厚爱健康管理有限公司",0),("让健康服务更可及，让本地消费更有活力",0)]),
]

base = '/home/ubuntu/.openclaw/workspace-jessica/'
build_deck(gov, base + '厚爱健康&生活大数据平台-政府汇报版.pptx', '政府汇报版')
build_deck(biz, base + '厚爱健康&生活大数据平台-招商版.pptx', '招商版')
build_deck(roadshow, base + '厚爱健康&生活大数据平台-综合路演版.pptx', '综合路演版')

print('done')
