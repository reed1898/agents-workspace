from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# colors
BLUE = RGBColor(20, 77, 163)
DARK = RGBColor(32, 33, 36)
GRAY = RGBColor(95, 99, 104)
LIGHT = RGBColor(245, 248, 255)


def set_bg(slide, color=RGBColor(255,255,255)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.9))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = DARK
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY


def add_section_bar(slide, text):
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(1.35), Inches(5.6), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = text
    tf.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(16)


def add_bullets(slide, items, left=0.9, top=2.0, width=11.5, height=4.8, level2_indent=22):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            text, level = it
        else:
            text, level = it, 0
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(20 if level==0 else 16)
        p.font.color.rgb = DARK if level==0 else GRAY
        if level == 0:
            p.font.bold = True
        p.space_after = Pt(8 if level==0 else 4)


def add_footer(slide, page):
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.01))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(220,225,235); line.line.fill.background()
    f = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(10), Inches(0.3)).text_frame
    f.text = "湖北厚爱健康管理有限公司｜业务方向综合统一版"
    f.paragraphs[0].font.size = Pt(10)
    f.paragraphs[0].font.color.rgb = GRAY
    pn = slide.shapes.add_textbox(Inches(12.2), Inches(6.86), Inches(0.5), Inches(0.3)).text_frame
    pn.text = str(page)
    pn.paragraphs[0].alignment = PP_ALIGN.RIGHT
    pn.paragraphs[0].font.size = Pt(10)
    pn.paragraphs[0].font.color.rgb = GRAY


# 1 cover
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT)
add_title(s, "厚爱健康业务方向综合方案", "整合版｜润色稿")
add_bullets(s, [
    "湖北厚爱健康管理有限公司",
    ("医疗信息化 · 社区健康服务 · 智慧养老",1),
    ("日期：2026.03",1)
], top=2.4)
add_footer(s, 1)

# 2 目录
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_title(s, "目录")
add_bullets(s, [
    "1. 公司定位与核心能力",
    "2. 四大业务方向总览",
    "3. 医疗信息化软件服务",
    "4. 厚爱健康APP与智能健康小包",
    "5. 社区健康小屋建设与运营",
    "6. 智慧养老解决方案",
    "7. 协同落地路径与下一步"
], top=1.8)
add_footer(s, 2)

# 3 定位
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_title(s, "公司定位与核心能力")
add_section_bar(s, "定位")
add_bullets(s, [
    "以“健康为本，爱满家园”为使命，打造覆盖全生命周期的健康服务生态",
    "形成“技术平台 + 场景服务 + 持续运营”的一体化能力",
    "聚焦三类核心客户：医疗机构、社区家庭、养老机构"
], top=2.0)
add_section_bar(s, "核心能力")
add_bullets(s, [
    "医疗信息化系统规划与交付",
    "健康数据采集、分析与闭环管理",
    "线上线下融合的社区与养老服务运营"
], top=4.35)
add_footer(s, 3)

# 4 总览
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_title(s, "四大业务方向总览")
add_bullets(s, [
    "医疗信息化软件服务：提升诊疗效率与数据治理能力",
    "厚爱健康APP + 智能健康小包：实现居家连续监测与慢病管理",
    "社区健康小屋：打通“家庭—社区—医院”服务闭环",
    "智慧养老解决方案：以智能化工具提升养老服务质量"
], top=2.1)
add_footer(s, 4)

# 5 医疗信息化
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_title(s, "业务方向一｜医疗信息化软件服务")
add_bullets(s, [
    "面向对象：医院、基层医疗机构、专科机构",
    "方案内容：系统集成、数据平台、流程数字化改造",
    "交付模式：需求调研 → 系统建设 → 培训上线 → 运维优化",
    "客户价值：效率提升、数据可视、决策更快、风险可控"
], top=2.0)
add_footer(s, 5)

# 6 APP+小包
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_title(s, "业务方向二｜厚爱健康APP与智能健康小包")
add_bullets(s, [
    "目标人群：中老年用户及其家庭成员",
    "核心功能：健康监测、风险预警、在线咨询、慢病随访",
    "设备协同：健康小包实现生理指标采集与APP实时同步",
    "产品优势：居家可用、家庭可见、服务可持续"
], top=2.0)
add_footer(s, 6)

# 7 社区小屋
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_title(s, "业务方向三｜社区健康小屋")
add_bullets(s, [
    "服务定位：社区端一站式健康触点",
    "核心服务：慢病筛查、健康咨询、中医诊疗、转诊协同",
    "运营机制：定时开放、专业医护、结果即时反馈",
    "社会价值：提升居民健康意识，增强基层服务可及性"
], top=2.0)
add_footer(s, 7)

# 8 合作运营
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_title(s, "社区场景合作与运营示例")
add_bullets(s, [
    "合作模式：联合地产/社区机构共建健康服务站点",
    "资源整合：场地资源 + 医疗能力 + 数字化运营",
    "运营目标：高频服务触达、持续健康管理、用户口碑沉淀",
    "可复制性：形成标准化建设与运营手册，支持快速复制"
], top=2.0)
add_footer(s, 8)

# 9 智慧养老
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_title(s, "业务方向四｜智慧养老解决方案")
add_bullets(s, [
    "机构端：智能巡检、物资管理、护理流程数字化",
    "家庭端：远程监护、异常预警、护理培训支持",
    "技术端：可穿戴设备 + 数据分析 + AI辅助决策",
    "体验端：兼顾安全、效率与老人生活质量"
], top=2.0)
add_footer(s, 9)

# 10 技术底座
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_title(s, "统一技术底座（支撑四大业务）")
add_bullets(s, [
    "云计算与数据中台：支撑多场景业务协同",
    "智能硬件接入能力：血压、心电等多设备兼容",
    "安全与合规：数据分级、权限控制、过程审计",
    "持续迭代机制：通过用户反馈驱动产品优化"
], top=2.0)
add_footer(s, 10)

# 11 落地路径
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
add_title(s, "落地路径与阶段目标")
add_bullets(s, [
    "阶段1（0-3个月）：完成重点场景试点与流程固化",
    "阶段2（3-6个月）：形成可复制方案并扩大覆盖范围",
    "阶段3（6-12个月）：建立跨场景数据协同与运营闭环",
    "关键抓手：示范案例、标准化手册、联合生态伙伴"
], top=2.0)
add_footer(s, 11)

# 12 ending
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, LIGHT)
add_title(s, "谢谢")
add_bullets(s, [
    "湖北厚爱健康管理有限公司",
    ("让健康服务更可及、更连续、更有温度。",1)
], top=2.8)
add_footer(s, 12)

out='/home/ubuntu/.openclaw/workspace-jessica/厚爱健康-业务方向综合统一-润色版.pptx'
prs.save(out)
print(out)
