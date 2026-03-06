from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()

BLUE = RGBColor(20, 77, 163)
DARK = RGBColor(32, 33, 36)
GRAY = RGBColor(95, 99, 104)
LIGHT = RGBColor(245, 248, 255)


def set_bg(slide, color=RGBColor(255,255,255)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.95))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = DARK
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY


def add_bullets(slide, items, left=0.85, top=1.7, width=11.6, height=4.9):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for text, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        p.font.size = Pt(18 if level == 0 else 15)
        p.font.color.rgb = DARK if level == 0 else GRAY
        p.font.bold = (level == 0)
        p.space_after = Pt(6 if level == 0 else 3)


def add_footer(slide, page):
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.01))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(220,225,235); line.line.fill.background()
    f = slide.shapes.add_textbox(Inches(0.6), Inches(6.9), Inches(10.5), Inches(0.3)).text_frame
    f.text = "湖北厚爱健康管理有限公司｜业务方向综合统一版（增强版）"
    f.paragraphs[0].font.size = Pt(10)
    f.paragraphs[0].font.color.rgb = GRAY
    pn = slide.shapes.add_textbox(Inches(12.2), Inches(6.86), Inches(0.5), Inches(0.3)).text_frame
    pn.text = str(page)
    pn.paragraphs[0].alignment = PP_ALIGN.RIGHT
    pn.paragraphs[0].font.size = Pt(10)
    pn.paragraphs[0].font.color.rgb = GRAY

slides = []

slides.append(("厚爱健康业务方向综合方案", "战略版｜用于对外汇报/合作沟通", [
    ("湖北厚爱健康管理有限公司",0),
    ("关键词：医疗信息化 · 社区健康服务 · 智慧养老",1),
    ("定位：成为区域健康服务数字化与运营一体化平台",1),
    ("日期：2026.03",1),
]))

slides.append(("目录", None, [
    ("1. 公司定位与核心竞争力",0),
    ("2. 市场机会与业务布局",0),
    ("3. 业务方向一：医疗信息化",0),
    ("4. 业务方向二：APP+智能健康小包",0),
    ("5. 业务方向三：社区健康小屋",0),
    ("6. 业务方向四：智慧养老",0),
    ("7. 商业模式与收入结构",0),
    ("8. 落地计划与合作诉求",0),
]))

slides.append(("公司定位与核心竞争力", None, [
    ("公司定位：健康服务“平台+场景+运营”综合服务商",0),
    ("围绕医院、社区、家庭、机构养老四类场景构建闭环",1),
    ("核心能力1：医疗信息化系统规划、建设与持续运维",0),
    ("核心能力2：健康数据采集—分析—干预—复盘闭环",0),
    ("核心能力3：线下服务运营与标准化复制能力",0),
    ("核心优势：既懂医疗业务流程，也懂数字化产品落地",0),
]))

slides.append(("市场机会与业务布局", None, [
    ("政策驱动：分级诊疗、医养结合、基层健康服务持续推进",0),
    ("需求驱动：慢病管理与居家健康监测需求高增长",0),
    ("公司布局：医疗机构端 + 居家端 + 社区端 + 养老端协同",0),
    ("战略目标：打造“高频服务入口 + 长周期健康管理”体系",0),
    ("增长逻辑：从单点项目收入升级为持续服务收入",0),
]))

slides.append(("业务方向一｜医疗信息化软件服务", None, [
    ("面向对象：医院、基层医疗机构、专科机构",0),
    ("核心模块：电子病历协同、运营数据看板、质控与随访",0),
    ("实施路径：调研诊断→方案设计→上线培训→运维优化",0),
    ("价值产出：提升诊疗效率、减少人工差错、增强管理透明",0),
    ("可交付成果：项目验收文档、培训体系、运维SLA机制",0),
]))

slides.append(("业务方向二｜厚爱健康APP与智能健康小包", None, [
    ("目标人群：慢病人群、中老年家庭、健康管理用户",0),
    ("核心功能：生理指标监测、风险预警、在线咨询、随访提醒",0),
    ("设备能力：血压/血糖/心率等设备接入与数据自动同步",0),
    ("服务价值：把“偶发就医”升级为“连续健康管理”",0),
    ("运营策略：用户分层、打卡激励、家庭成员协同管理",0),
]))

slides.append(("业务方向三｜社区健康小屋", None, [
    ("定位：社区居民可及的一站式健康服务触点",0),
    ("服务组合：筛查评估、健康咨询、慢病跟踪、转诊协同",0),
    ("运营机制：固定时段+主题活动+数据回访",0),
    ("合作模式：街道/物业/社区机构共建共营",0),
    ("预期成效：提升基层可及性，形成社区健康品牌影响力",0),
]))

slides.append(("业务方向四｜智慧养老解决方案", None, [
    ("机构端：护理流程数字化、风险巡检、物资与工单管理",0),
    ("家庭端：远程看护、异常告警、照护指导",0),
    ("技术端：可穿戴设备接入 + 数据分析 + AI辅助评估",0),
    ("管理端：运营报表、服务质量追踪、满意度闭环",0),
    ("核心价值：兼顾安全、效率与长者生活质量",0),
]))

slides.append(("商业模式与收入结构", None, [
    ("收入来源A：项目建设与系统交付收入",0),
    ("收入来源B：平台订阅与运维服务年费",0),
    ("收入来源C：健康服务包与增值服务收入",0),
    ("收入来源D：社区站点与机构端合作分成",0),
    ("目标结构：提高经常性收入占比，增强现金流稳定性",0),
]))

slides.append(("统一技术底座与数据治理", None, [
    ("技术架构：云平台+数据中台+多终端应用协同",0),
    ("数据能力：统一身份、统一标签、统一服务编排",0),
    ("设备兼容：支持多品牌健康硬件标准化接入",0),
    ("安全合规：权限分级、审计追踪、隐私保护机制",0),
    ("迭代机制：以用户反馈和运营数据驱动产品优化",0),
]))

slides.append(("阶段化落地计划（12个月）", None, [
    ("0-3个月：试点项目上线，打通核心流程与数据链路",0),
    ("3-6个月：沉淀标准化方案，复制到更多场景与区域",0),
    ("6-12个月：形成跨场景协同与持续运营闭环",0),
    ("关键保障：联合生态伙伴、建立示范案例、强化服务SOP",0),
    ("建议KPI：服务覆盖人数、活跃率、复购率、满意度",0),
]))

slides.append(("合作诉求与下一步", None, [
    ("合作方向：渠道合作、场景共建、技术协同、资本合作",0),
    ("近期目标：落地样板点，形成可展示可复制案例",0),
    ("沟通机制：周节奏推进 + 月度复盘 + 季度里程碑",0),
    ("期待成果：共同打造区域健康服务标杆项目",0),
    ("谢谢｜让健康服务更可及、更连续、更有温度",0),
]))

for i, (title, subtitle, items) in enumerate(slides, 1):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, LIGHT if i in [1,12] else RGBColor(255,255,255))
    add_title(s, title, subtitle)
    add_bullets(s, items)
    add_footer(s, i)

out = '/home/ubuntu/.openclaw/workspace-jessica/厚爱健康-业务方向综合统一-增强版.pptx'
prs.save(out)
print(out)
